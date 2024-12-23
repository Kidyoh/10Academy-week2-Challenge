"""Presentation generator for interim submission."""
import pptx
from pptx.util import Inches, Pt
from pathlib import Path
import logging
from src.utils.visualization_helper import save_visualizations

logger = logging.getLogger('tellco_analysis')

class PresentationGenerator:
    """Class for generating presentation slides."""
    
    def __init__(self, output_dir="presentations"):
        self.prs = pptx.Presentation()
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(exist_ok=True)
        logger.info(f"Initialized presentation generator. Output dir: {self.output_dir}")
    
    def add_title_slide(self, title, subtitle):
        """Add title slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[0])
        slide.shapes.title.text = title
        slide.placeholders[1].text = subtitle
        logger.info("Added title slide")
    
    def add_section_header(self, title):
        """Add section header slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[2])
        slide.shapes.title.text = title
        logger.info(f"Added section header: {title}")
    
    def add_bullet_slide(self, title, bullets):
        """Add slide with bullet points."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        slide.shapes.title.text = title
        
        # Add bullet points
        body_shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame
        
        for bullet in bullets:
            p = tf.add_paragraph()
            p.text = bullet
            p.level = 0
        
        logger.info(f"Added bullet slide: {title}")
    
    def add_image_slide(self, title: str, image_path: str, content: list):
        """Add slide with image and bullet points."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        slide.shapes.title.text = title
        
        # Add image
        left = Inches(1)
        top = Inches(2)
        width = Inches(4)
        slide.shapes.add_picture(image_path, left, top, width=width)
        
        # Add bullet points
        text_box = slide.shapes.add_textbox(Inches(5.5), Inches(2), Inches(4), Inches(5))
        tf = text_box.text_frame
        
        for bullet in content:
            p = tf.add_paragraph()
            p.text = bullet
            p.level = 0
    
    def generate_presentation(self, task1_results, engagement_results):
        """Generate complete presentation with visualizations."""
        try:
            # Save visualizations
            viz_paths = save_visualizations({**task1_results, **engagement_results})
            
            # Title Slide
            self.add_title_slide(
                "TellCo Telecom Analysis",
                "Interim Analysis Results"
            )
            
            # Task 1 Section
            self.add_section_header("Task 1: User Overview Analysis")
            
            # Handset Analysis with Visualization
            self.add_image_slide(
                "Top 10 Handsets",
                viz_paths['top_handsets'],
                [
                    f"Most popular: {task1_results['top_handsets'].iloc[0]['handset']}",
                    f"Second: {task1_results['top_handsets'].iloc[1]['handset']}",
                    "Apple devices dominate the top 10",
                    "High-end smartphones preferred"
                ]
            )
            
            # Manufacturer Analysis with Visualization
            self.add_image_slide(
                "Manufacturer Market Share",
                viz_paths['manufacturer_share'],
                [
                    f"Leader: {task1_results['top_manufacturers'].iloc[0]['manufacturer']}",
                    "Market dominated by three manufacturers",
                    "Opportunity for diversification",
                    "Premium brand preference"
                ]
            )
            
            # Task 2 Section
            self.add_section_header("Task 2: User Engagement Analysis")
            
            # Engagement Patterns with Visualization
            self.add_image_slide(
                "User Engagement Patterns",
                viz_paths['engagement_pattern'],
                [
                    "Clear correlation between duration and traffic",
                    "High-end device users show longer sessions",
                    "Distinct usage patterns identified",
                    "Potential for targeted optimization"
                ]
            )
            
            # Application Usage with Visualization
            self.add_image_slide(
                "Application Usage Distribution",
                viz_paths['app_usage'],
                [
                    "Social media dominates traffic",
                    "Video streaming significant",
                    "Gaming shows growth potential",
                    "Email remains steady"
                ]
            )
            
            # Save presentation
            output_path = self.output_dir / "TellCo_Analysis_Interim.pptx"
            self.prs.save(str(output_path))
            
        except Exception as e:
            logger.error(f"Error generating presentation: {str(e)}")
            raise
    
    def save(self, filename):
        """Save the presentation."""
        output_path = self.output_dir / filename
        try:
            self.prs.save(str(output_path))
            logger.info(f"Presentation saved to: {output_path}")
        except Exception as e:
            logger.error(f"Error saving presentation: {str(e)}")
            raise